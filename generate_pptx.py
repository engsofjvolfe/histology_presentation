# gerar_pptx.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

prs = Presentation()
prs.slide_width = Inches(13.33)  # widescreen
prs.slide_height = Inches(7.5)

def add_slide(title, bullets, footer_text=None, img_placeholder_text=None, notes_text=None):
    slide_layout = prs.slide_layouts[1]  # title and content
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    title_tf.text = title

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i>0 else body.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(18)

    # footer (small)
    if footer_text:
        tx = slide.shapes.add_textbox(Inches(0.2), Inches(7.0), Inches(12.9), Inches(0.4))
        p = tx.text_frame.paragraphs[0]
        p.text = footer_text
        p.font.size = Pt(10)

    # image placeholder box
    if img_placeholder_text:
        left = Inches(9.3)
        top = Inches(1.2)
        width = Inches(3.4)
        height = Inches(3.6)
        ph = slide.shapes.add_textbox(left, top, width, height)
        p = ph.text_frame.paragraphs[0]
        p.text = img_placeholder_text
        p.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p.font.size = Pt(12)

    # speaker notes
    if notes_text:
        notes = slide.notes_slide.notes_text_frame
        notes.text = notes_text

# ========== Slides content (25 slides) ==========
slides = [
    # 1
    ("Conectando Remodelação e Regeneração Óssea",
     ["Desvendando hormônios e vias de sinalização", "Seminário de Histologia", "Mehreen et al., Biology (2025) — DOI: 10.3390/biology14030274"],
     "Slide 1 • PDF: p.1"),
    # 2
    ("Equipe de Apresentação",
     ["Apresentador 01: Introdução e Desafios",
      "Apresentador 02: Organismos Modelo",
      "Apresentador 03: Fases da Regeneração",
      "Apresentador 04: Remodelação e PTH",
      "Apresentador 05: Calcitonina e FGF23",
      "Apresentador 06: IGF/GH e Conclusões"],
     "Slide 2 • Institucional"),
    # 3
    ("Agenda",
     ["Introdução e Contexto — Ap.01",
      "Organismos Modelo — Ap.02",
      "Fases da Regeneração — Ap.03",
      "Remodelação e PTH — Ap.04",
      "Calcitonina e FGF23 — Ap.05",
      "IGF/GH e Conclusões — Ap.06"],
     "Slide 3 • PDF: p.2"),
    # 4
    ("Introdução: Importância Clínica",
     ["Milhões afetados por doenças ósseas", "Defeitos críticos (>2 cm) não cicatrizam espontaneamente", "Soluções: engenharia tecidual, modulação hormonal, biomateriais"],
     "Apresentador 01 | PDF: p.1-2"),
    # 5
    ("Objetivos e Abordagem",
     ["Objetivo: conectar regeneração (modelos) ↔ remodelação (humanos)", "Foco: Wnt/BMP/TGF-β e regulação hormonal", "Avaliar riscos e benefícios de manipulação hormonal"],
     "Apresentador 01 | PDF: p.2"),
    # 6
    ("Desafios Translacionais",
     ["Proliferação descontrolada e risco oncogênico", "Crosstalk complexo entre vias (Wnt/BMP/TGF-β)", "Necessidade de entrega localizada e temporização precisa"],
     "Apresentador 01 | PDF: p.2-3"),
    # 7
    ("Organismos Modelo — Axolote & Zebrafish",
     ["Axolote: regeneração completa, blastema verdadeiro, neotenia", "Zebrafish: manipulação genética, nadadeira ≈ osso, vias conservadas", "T3/T4 e PTHrP estão implicados nas respostas regenerativas"],
     "Apresentador 02 | PDF: p.3-4",
     "[FIGURA: inserir foto Axolote / Zebrafish]"),
    # 8
    ("Organismos Modelo — Xenopus e Comparações",
     ["Xenopus: girino regenerativo, adulto perde capacidade", "Pico de T3 na metamorfose reduz regeneração", "Comparativo: anfíbios vs mamíferos (limitações)"],
     "Apresentador 02 | PDF: p.4-5"),
    # 9
    ("Regulação Hormonal (Tabela 1) — Parte 1",
     ["T3/T4: metabolismo e crescimento (papéis dependentes de espécie)","IGF-1: proliferação de progenitoras","BMPs: desdiferenciação e osteogênese"],
     "Apresentador 02 | PDF: p.5 (Tabela 1)"),
    # 10
    ("Regulação Hormonal — Parte 2",
     ["VEGF: angiogênese essencial","Ácido Retinoico: efeitos dualistas dose-dependentes","GH: suporte via IGF-1, sinergismo com fatores de crescimento"],
     "Apresentador 02 | PDF: p.5-6"),
    # 11
    ("Fase Inicial: Hemostasia e Inflamação",
     ["Neutrófilos → macrófagos; liberação de citocinas","Recrutamento de progenitoras locais","Ativação hormonal (PTHrP, Sox9) influencia condrogênese)"],
     "Apresentador 03 | PDF: p.6-7"),
    # 12
    ("Regeneração Epimórfica & Blastema",
     ["Wound epidermis → proteção + sinalização","Inervação necessária para blastema","Desdiferenciação → massa proliferativa (blastema)"],
     "Apresentador 03 | PDF: p.7 (Figura 1)",
     "[FIGURA 1: inserir imagem do PDF — Axolotl limb regeneration]"),
    # 13
    ("Desdiferenciação e Regulação Molecular",
     ["Osteoblastos → progenitoras multipotentes (BMPs, T3)","BMP2b/BMP6 e Shh → iniciam proliferação","Fatores: Msx1 (indiferenciação), Pax7 (proliferação)"],
     "Apresentador 03 | PDF: p.7-8"),
    # 14
    ("Proliferação e Diferenciação",
     ["IGF-1: reentrada no ciclo celular","Shift metabólico: glicólise ↑, OXPHOS ↓ (Zebrafish)","BMP → Runx2 → diferenciação osteoblástica e mineralização"],
     "Apresentador 03 | PDF: p.8"),
    # 15
    ("Remodelação Óssea — Ciclo (humanos)",
     ["Ativação (estresse mecânico / PTH)","Reabsorção (osteoclastos, RANKL)","Reversão e Formação (osteoblastos, mineralização)"],
     "Apresentador 04 | PDF: p.8-9",
     "[FIGURA 2: inserir imagem do PDF — Bone Remodeling]"),
    # 16
    ("Regulação Hormonal da Remodelação",
     ["RANKL/OPG regula osteoclastogênese","Calcitonina inibe osteoclastos","Wnt/β-catenina e RUNX2 promovem formação óssea"],
     "Apresentador 04 | PDF: p.9 (Figura 2)"),
    # 17
    ("PTH — Mecanismos Moleculares",
     ["PTH → PTH1R → cAMP/PKA","Inativação GSK-3β → estabiliza β-catenina","Ativação de genes osteogénicos (Runx2, osteocalcina)"],
     "Apresentador 04 | PDF: p.10"),
    # 18
    ("PTH — Regulação Dual e Aplicações",
     ["Pulsátil → formação óssea anabólica","Contínuo → aumento da reabsorção","Teriparatida: uso em osteoporose e não-uniões; local pulsátil em scaffolds"],
     "Apresentador 04 | PDF: p.10-11"),
    # 19
    ("Calcitonina — Funções principais",
     ["Hipocalcemiante; inibe osteoclastos","Estabiliza β-catenina → suporte à osteogênese","Regula proliferação do blastema (evita hiperproliferação)"],
     "Apresentador 05 | PDF: p.11-12"),
    # 20
    ("FGF23 — Fosfato e Mineralização",
     ["Regula reabsorção renal de fosfato","Suprime síntese de calcitriol","Níveis elevados → hipofosfatemia e prejuízo da mineralização"],
     "Apresentador 05 | PDF: p.11-13"),
    # 21
    ("Integração Calcitonina ↔ FGF23",
     ["Cálcio (calcitonina) vs Fosfato (FGF23) — equilíbrio essencial","Desbalanços → osteomalacia/raquitismo","Abordagens: liberação localizada, scaffolds biomiméticos"],
     "Apresentador 05 | PDF: p.12-14"),
    # 22
    ("IGF-1 & GH — Papel e Riscos",
     ["IGF-1: proliferação e diferenciação osteoblástica (PI3K/Akt, MAPK)","GH: aumenta IGF-1 hepático/local","Riscos: oncogênese, senescência, efeitos metabólicos"],
     "Apresentador 06 | PDF: p.14-15"),
    # 23
    ("Interações e Perspectivas Futuras",
     ["IGF/GH ↔ BMP/TGF-β/Wnt — decisões osteogênese vs fibrose","Biomateriais sensíveis a hormônios","Medicina de precisão: entregas locais e temporizadas"],
     "Apresentador 06 | PDF: p.15-16"),
    # 24
    ("Conclusões",
     ["Vias conservadas: Wnt, BMP, TGF-β","Hormônios centrais: PTH, Calcitonina, FGF23, IGF/GH","Tradução clínica: equilíbrio eficácia × segurança; intervenção local"],
     "Apresentador 06 | PDF: p.15-17"),
    # 25
    ("Agradecimentos e Perguntas",
     ["Autores: Mehreen et al., Biology (2025)","Licença: CC BY 4.0 — DOI: 10.3390/biology14030274","Orientadores e instituição"],
     "Slide final • PDF: p.15-17",
     "[FIGURA: esquema final; inserir imagem]"),
]

for s in slides:
    # s may have 3 or 4 or 5 items
    title = s[0]
    bullets = s[1]
    footer = s[2] if len(s) > 2 else None
    img_ph = s[3] if len(s) > 3 else None
    notes = s[4] if len(s) > 4 else None
    add_slide(title, bullets, footer_text=footer, img_placeholder_text=img_ph, notes_text=notes)

# Save file
prs.save("presentation_final.pptx")
print("presentation_final.pptx criado com sucesso.")
